import os
import yaml
from openai import OpenAI

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import requests
from PIL import Image

from datetime import datetime 

OPENAI_API_KEY = os.getenv('aitui_openai_key')
client = OpenAI(api_key=OPENAI_API_KEY)

blue_color = RGBColor(50, 40, 255)
black_color = RGBColor(22, 12, 32)
grey_color = RGBColor(169, 169, 169)

def create_dalle_image(image_desc, filename):
    print('\nImage Prompt:', image_desc, '\nTitle:', filename)
    
    response = client.images.generate(
        model="dall-e-3",
        prompt=image_desc,
        size="1024x1024",
        quality="standard",
        n=1,
        )

    image_url = response.data[0].url
    image_data = requests.get(image_url).content
    image_filename = filename
    image_path = image_filename
    if 'assets' not in filename:
        image_path = os.path.join("assets",image_filename)
    with open(image_path, "wb") as f:
        f.write(image_data)
    return image_path
def add_intro_slide(prs, data):
    intro_slide_layout = prs.slide_layouts[1]
    intro_slide = prs.slides.add_slide(intro_slide_layout)

    # Add author and date at the bottom left
    left = Inches(0.25)
    top = Inches(6.5)
    width = Inches(8)
    height = Inches(0.5)
    author_date = intro_slide.shapes.add_textbox(left, top, width, height)
    author_date_text_frame = author_date.text_frame
    author_date_text_frame.word_wrap = True
    author_date_text_frame.text = f"Author: {data['author']}\nCreated Date: {data['date']}"
    for paragraph in author_date_text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(16)

    # Add title top left
    title = intro_slide.shapes.title
    title.text = data['title']
    title.text_frame.paragraphs[0].font.color.rgb = blue_color
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Generate image for intro slide
    image_url=data.get("intro_image")
    pic = intro_slide.shapes.add_picture(image_url, Inches(4), Inches(0.5), height=Inches(6.5))
    
 
    pic.left = Inches(5)
    pic.top = Inches(00)
    pic.width = Inches(5)
    pic.height = Inches(7.5)

def add_exit_slide(prs, data):
    exit_slide_layout = prs.slide_layouts[5] # Blank layout for a final slide
    exit_slide = prs.slides.add_slide(exit_slide_layout)
    
    # Open the image and get its dimensions
    with Image.open(data['logo_path']) as img:
        img_width, img_height = img.size
    
    # Calculate aspect ratio
    aspect_ratio = img_width / img_height
    
    # Define desired dimensions for the image
    desired_width = Inches(5)  # Adjust as needed
    desired_height = desired_width / aspect_ratio
    
    # Calculate top position to center the image vertically
    top = (prs.slide_height - desired_height) / 2
    
    # Add picture with correct size and aspect ratio
    exit_slide.shapes.add_picture(data['logo_path'], Inches(2.5), top, width=desired_width, height=desired_height)
def create_ai_adoption_ppt(template_file):
    # Load data from YAML template
    with open(template_file, 'r') as file:
        data = yaml.safe_load(file)

    # Create a PowerPoint presentation
    prs = Presentation()

    # Define IBM blue color
    ibm_blue = RGBColor(0, 70, 134)

    # Set IBM Plex font and adjust slide font size
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = "IBM Plex Sans"
                        run.font.size = Pt(24)
        for placeholder in slide.placeholders:
            for paragraph in placeholder.text_frame.paragraphs:
                paragraph.font.size = Pt(20)

 
    # Add intro slide
    add_intro_slide(prs, data)

    # Add sections slides with detailed content and blue gradient background
    for section in data['sections']:
        slide_layout = prs.slide_layouts[7]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = section["title"]
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = blue_color

        content = slide.placeholders[1]
        content.text = section["content"]
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(18)

        content.text_frame.paragraphs[0].font.size = Pt(18)
        content.text_frame.paragraphs[0].font.color.rgb = black_color

        # Apply professional template background (e.g., high-quality image)
        left = Inches(0) 
        top = Inches(0)
        pic = slide.shapes.add_picture(data.get("background_image"), left, top, width=prs.slide_width, height=prs.slide_height)
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)  # Adjust index as needed to send it to the back

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = grey_color
        slide.background.fill.transparency = 0.2


        image_path = section.get('image')
        if image_path and os.path.exists(image_path)==False:
            image_path =create_dalle_image(section["title"],image_path)  
        left = Inches(.5) 
        top = Inches(1.57)
        
        pic = slide.shapes.add_picture(image_path,left,top, width=Inches(3), height=Inches(3))


    # Add exit slide
    add_exit_slide(prs, data)

    # Save the presentation
    output_file = data.get('ppt_name').replace(' ','_')
    prs.save(output_file)
    print(f"Presentation generated: {output_file}")

if __name__ == "__main__":
    create_ai_adoption_ppt("template.yaml")
